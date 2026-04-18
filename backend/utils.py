import os
import re
import PyPDF2
import io

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def extract_text_from_pdf(file_stream):
    text_content = []

    try:
        reader = PyPDF2.PdfReader(file_stream)

        for page_num, page in enumerate(reader.pages):
            text = page.extract_text()

            if text:
                text = text.strip()
                if text:
                    text_content.append({
                        "page": page_num + 1,
                        "text": text
                    })

    except Exception as e:
        print(f"PDF extraction error: {e}")

    return text_content


def extract_text_from_docx(file_stream):
    if docx is None:
        return []

    try:
        document = docx.Document(file_stream)

        text = "\n".join(
            [p.text for p in document.paragraphs if p.text and p.text.strip()]
        )

        if text:
            return [{"page": 1, "text": text.strip()}]

    except Exception as e:
        print(f"DOCX error: {e}")

    return []


def extract_text_from_pptx(file_stream):
    if Presentation is None:
        return []

    try:
        prs = Presentation(file_stream)
        slides_text = []

        for slide in prs.slides:
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    t = shape.text.strip()
                    if t:
                        slide_text.append(t)

            if slide_text:
                slides_text.append("\n".join(slide_text))

        if slides_text:
            return [{"page": 1, "text": "\n\n".join(slides_text)}]

    except Exception as e:
        print(f"PPTX error: {e}")

    return []


def extract_text_from_file(file_path, file_stream=None):
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    if file_stream is None:
        try:
            with open(file_path, "rb") as f:
                file_stream = io.BytesIO(f.read())
        except Exception:
            return []
    
    file_stream.seek(0)

    if ext == ".pdf":
        return extract_text_from_pdf(file_stream)
    if ext == ".docx":
        return extract_text_from_docx(file_stream)
    if ext == ".pptx":
        return extract_text_from_pptx(file_stream)

    if ext in [".txt", ".md", ".csv", ".json", ".log", ".html", ".xml"]:
        try:
            text = file_stream.read().decode("utf-8", errors="ignore")
            text = text.strip()
            if text:
                return [{"page": 1, "text": text}]

        except Exception as e:
            print(f"Text file error: {e}")

    return []


def chunk_text(text_data, chunk_size=900, overlap=150, max_chunks=None):
    if not text_data or not isinstance(text_data, list):
        return []

    if max_chunks is not None:
        try:
            max_chunks = int(max_chunks)
        except Exception:
            max_chunks = None

    if max_chunks is not None and max_chunks <= 0:
        return []

    chunks = []
    step = max(1, int(chunk_size) - int(overlap))

    for item in text_data:
        if not isinstance(item, dict):
            continue

        text = str(item.get("text") or "")
        page = item.get("page", 1)

        text = re.sub(r"\s+", " ", text).strip()
        if not text:
            continue

        start = 0
        while start < len(text):
            end = start + int(chunk_size)
            chunk = text[start:end].strip()

            if len(chunk) > 50:
                chunks.append(
                    {
                        "content": chunk,
                        "metadata": {"page": page},
                    }
                )
                if max_chunks is not None and len(chunks) >= max_chunks:
                    return chunks

            start += step

    return chunks
